from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

class SlideAdder():
    def __init__(self, presentation, config:dict={}):
        """
        Single slide parser

        Arguments:
         - presentation :pptx.Presentation for slide adding
         - config       :dict containing data about font and layout number
        """
        self.prs = presentation

        self.font_name = config.get("font_name", "Arial")
        self.font_size = config.get("font_size",      16)
        self.slide_layout = self.prs.slide_layouts[config.get("layout", 5)]

    def __call__(self, slide_config:dict):
        """
        Arguments:
         - slide_config : dict
        """
        text_width = width / (2 if slide_config["figures"] else 1) - 0.5

        slide = self.prs.slides.add_slide(self.slide_layout)
        
        title = slide.shapes.title
        self.title(title, slide_config["title"])
        
        idea_frame = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(width - 1), Inches(0.5)
        )
        self.idea_subtitle(idea_frame, slide_config["idea"])

        text_frame = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(text_width), Inches(4)
        )
        self.text_bullets(text_frame, slide_config["text"])

        y_offset = Inches(2.5) + Inches(0.4) * len(text_frame.paragraphs)
        add_latex_formulas_as_images(
            slide=slide, 
            slide_data=slide_config.get("formulas", None), 
            y_offset=y_offset, 
            text_width=text_width,
        )

    def text_bullets(self, textbox, text):
        text_frame = textbox.text_frame
        text_frame.word_wrap = True # Enable word wrap
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.paragraphs[0].font.name = self.font_name

        buf = ""
        first = True
        for item in text:
            if buf.startswith("• ") and not item.startswith("• "):
                text_frame.add_paragraph()
            
            p = text_frame.add_paragraph() if not first else text_frame.paragraphs[0]
            p.font.name = self.font_name
            p.font.size = Pt(self.font_size)
            
            if item.startswith("• "):
                p.text = item[2:]  # Remove the "• " prefix
                p.level = 1
            elif item.startswith("- "):
                p.text = item[2:]  # Remove the "- " prefix
                p.level = 0
            else:
                p.text = item
                p.level = 0
                p.font.bold = True
            buf = item
            first = False

    def idea_subtitle(self, subtitle, text):
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = text
        subtitle_frame.paragraphs[0].font.size = Pt(self.font_size)
        subtitle_frame.paragraphs[0].font.name = self.font_name

    def title(self, title, text):
        text = text if text else "Test slide"
        title.text = text
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title.text_frame.word_wrap = False
        title.text_frame.paragraphs[0].font.name = self.font_name