import json
from os import PathLike, makedirs
from os.path import dirname
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from src.tools import add_latex_formulas_as_images

def create_presentation_from_json(
    json_file:PathLike, 
    output_pptx:PathLike,
    height:float=7.5,
    width:float=13.333,
):
    """
    Creates a PowerPoint presentation from a JSON file.

    Args:
        json_file: Path to the JSON file.
        output_pptx: Path to save the generated PowerPoint presentation.
    """
    FONT_SIZE = 16
    with open(json_file, 'r') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    # Use a slide layout with a title and content placeholders
    # Choose a layout that best suits your needs
    # Layout 0 is typically a title slide, and layout 1 is title and content
    slide_layout = prs.slide_layouts[5]

    for slide_key, slide_data in data.items():
        if slide_key.startswith("Slide"):
            slide = prs.slides.add_slide(slide_layout)
            text_width = width / (2 if slide_data["figures"] else 1) - 0.5
            
            # Title
            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title.text_frame.word_wrap = False
            title.text_frame.paragraphs[0].font.name = "Arial"

            # [print(f"{i} : {x}") for i, x in enumerate(dir(title.text_frame))]

            # Idea (smaller font underneath the title)
            idea_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(width - 1), Inches(0.5))
            idea_text_frame = idea_textbox.text_frame
            idea_text_frame.text = slide_data["idea"]
            idea_text_frame.paragraphs[0].font.size = Pt(FONT_SIZE)  # Adjust font size as needed
            idea_text_frame.paragraphs[0].font.name = "Arial"
            
            # Text (bullet list on the left half of the slide)
            text_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(text_width), Inches(4))
            text_frame = text_textbox.text_frame
            text_frame.word_wrap = True # Enable word wrap
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            text_frame.paragraphs[0].font.name = "Arial"

            buf = ""
            first = True
            for item in slide_data["text"]:
                if buf.startswith("• ") and not item.startswith("• "):
                    text_frame.add_paragraph()
                
                p = text_frame.add_paragraph() if not first else text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(FONT_SIZE)
                
                if item.startswith("• "):
                    p.text = item[2:]  # Remove the "• " prefix
                    p.level = 1
                elif item.startswith("- "):
                    p.text = item[2:]  # Remove the "- " prefix
                    p.level = 0
                else:
                    p.text = item
                    p.level = 0
                    p.font.bold = True
                buf = item
                first = False
                    
                
            y_offset = Inches(2.5) + Inches(0.4) * len(text_frame.paragraphs)

            add_latex_formulas_as_images(
                slide=slide, 
                slide_data=slide_data, 
                y_offset=y_offset, 
                text_width=text_width,
            )

            # Figures (placeholder image)
            figsize = {
                "left" : Inches(width / 2),
                "top" : Inches(2),
                "width" : Inches(width / 2 - 1),
                "height" : Inches(4),
            }
            
            # Add a placeholder picture. You can use a dummy image or a specific placeholder image
            # Replace 'path/to/placeholder.png' with an actual image file if you have one
            if slide_data["figures"]:
                try:
                    slide.shapes.add_picture('./placeholder/placeholder.png', figsize["left"], figsize["top"], figsize["width"], figsize["height"])
                except FileNotFoundError:
                    print("Warning: Placeholder image not found. Using a text placeholder instead.")
                    picture_placeholder = slide.shapes.add_textbox(figsize["left"], figsize["top"], figsize["width"], figsize["height"])
                    picture_placeholder.text_frame.text = f"Figure: {slide_data['figures']}"
                    picture_placeholder.text_frame.paragraphs[0].font.size = Pt(FONT_SIZE)

            # Speech (notes)
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speech"]

    makedirs(dirname(output_pptx), exist_ok=True)
    prs.save(output_pptx)

def generate_placeholder():
    ph = Image.new("RGB", (100, 100))
    makedirs("placeholder", exist_ok=True)    
    ph.save("placeholder/placeholder.png")

from argparse import ArgumentParser
if __name__ == "__main__":
    p = ArgumentParser()
    p.add_argument("--test"     , type=bool  , default=0                                 , help="If not 0 - creates dummy JSON for the test" )
    p.add_argument("--json"     , type=str   , default="slide_json/slides_data.json"     , help="Path to JSON with slides data"              )
    p.add_argument("--savepath" , type=str   , default="output/_output_presentation.pptx", help="Path for .PPTX to be saved"                 )
    args = p.parse_args()
    
    test = args.test
    json_file_path = args.json  # Replace with your JSON file path
    output_pptx_path = args.savepath  # Replace with desired output path


    if test:
        # Example usage:
        # Create a dummy JSON file for testing
        dummy_data = {
            "Slide 1": {
                "title": "Sample Title 1",
                "idea": "Sample Idea 1",
                "text": [
                    "- bullet1",
                    "- bullet2",
                ],
                "formulas": [r"\int_0^1 x^2 dx"],
                "figures": "Figure Description 1",
                "speech": "Additional Notes 1"
            },
            "Slide 2": {
                "title": "Sample Title 2",
                "idea": "Sample Idea 2",
                "text": [
                    "subtitle1:",
                    "• bullet1",
                    "• bullet2",
                    "subtitle2:",
                    "• bullet3",
                    "• bullet4",
                ],
                "formulas": [r"\sum_{i=1}^n i = \frac{n(n+1)}{2}"],
                "figures": "Figure Description 2",
                "speech": "Additional Notes 2"
            }
        }

        with open(json_file_path, 'w') as f:
            json.dump(dummy_data, f, indent=4)

    generate_placeholder()
    create_presentation_from_json(json_file_path, output_pptx_path)